#!/usr/bin/env python3
"""Create the balls-into-bins technical interview presentation."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "presentation" / "balls_into_bins_presentation.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

COLORS = {
    "ink": RGBColor(30, 41, 59),
    "muted": RGBColor(71, 85, 105),
    "light": RGBColor(248, 250, 252),
    "line": RGBColor(203, 213, 225),
    "blue": RGBColor(37, 99, 235),
    "green": RGBColor(22, 163, 74),
    "orange": RGBColor(234, 88, 12),
    "gray": RGBColor(100, 116, 139),
    "white": RGBColor(255, 255, 255),
}


def add_textbox(slide, x, y, w, h, text, font_size=22, bold=False,
                color="ink", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(0.05)
    box.text_frame.margin_right = Inches(0.05)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.45),
                title, 27, True)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.78), Inches(11.8), Inches(0.35),
                    subtitle, 13, False, "muted")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.18),
                                  Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.color.rgb = COLORS["line"]


def add_takeaway(slide, text, x=0.75, y=6.35, w=11.85, h=0.65, color="blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(239, 246, 255)
    shape.line.color.rgb = COLORS[color]
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = COLORS["ink"]


def add_bullets(slide, x, y, w, h, bullets, font_size=18, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(8)
    return box


def add_callout(slide, x, y, w, h, lines, color="blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS[color]
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(12.5 if len(line) > 45 else 14)
        p.font.bold = idx == 0
        p.font.color.rgb = COLORS["ink"] if idx == 0 else COLORS["muted"]
        p.space_after = Pt(4)
    return shape


def add_plot(slide, rel_path, x=0.72, y=1.35, w=8.45, h=4.75):
    path = ROOT / rel_path
    if not path.exists():
        add_placeholder(slide, x, y, w, h, f"Missing plot:\n{rel_path}")
        return False
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w),
                             height=Inches(h))
    return True


def add_placeholder(slide, x, y, w, h, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["light"]
    shape.line.color.rgb = COLORS["line"]
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(18)
    r.font.color.rgb = COLORS["gray"]


def add_table(slide, x, y, w, h, rows):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x),
                                         Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for col in range(len(rows[0])):
        table.columns[col].width = Inches(w / len(rows[0]))
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(226, 232, 240) if r_idx == 0 else COLORS["white"]
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(12 if len(value) > 24 else 14)
                p.font.bold = r_idx == 0
                p.font.color.rgb = COLORS["ink"]
    return table_shape


def add_pipeline(slide):
    labels = ["Scenarios", "Simulators", "Metrics CSV", "Plots", "Conclusions"]
    xs = [0.9, 3.0, 5.2, 7.45, 9.65]
    for idx, (label, x) in enumerate(zip(labels, xs)):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                       Inches(3.0), Inches(1.55), Inches(0.72))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["white"]
        shape.line.color.rgb = COLORS["blue"] if idx < 3 else COLORS["orange"]
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Aptos"
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = COLORS["ink"]
        if idx < len(labels) - 1:
            add_textbox(slide, Inches(x + 1.63), Inches(3.17), Inches(0.55),
                        Inches(0.25), "->", 20, True, "gray", PP_ALIGN.CENTER)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["light"]
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = blank_slide(prs)
    add_textbox(slide, Inches(0.75), Inches(1.0), Inches(11.7), Inches(0.8),
                "Balls into Bins: Load Balancing Under Cost Constraints",
                34, True)
    add_textbox(slide, Inches(0.78), Inches(1.9), Inches(10.8), Inches(0.45),
                "Randomized choices, stateful policies, and cost-quality tradeoffs",
                19, False, "muted")
    add_callout(slide, 0.85, 4.45, 6.5, 1.1,
                ["Main message",
                 "How much load-balance quality does each policy buy for its decision cost?"],
                "blue")
    add_bullets(slide, 8.0, 4.0, 4.0, 1.6,
                ["10,000 balls", "16 and 32 bins", "10 trials per scenario"], 18)

    slide = blank_slide(prs)
    add_title(slide, "Problem Statement")
    add_bullets(slide, 0.8, 1.55, 5.45, 2.5,
                ["Allocate each arriving ball immediately.",
                 "Goal: reduce max load and imbalance.",
                 "Systems analogy: assign tasks, flows, or jobs to resources.",
                 "Question: how much state or load information should the allocator use?"], 18)
    rows = [["Workload", "Initial bins"],
            ["unweighted", "empty"],
            ["weighted", "empty"],
            ["unweighted", "randomly initialized"],
            ["weighted", "randomly initialized"]]
    add_table(slide, 7.0, 1.55, 4.9, 2.65, rows)
    add_takeaway(slide, "The experiments compare online allocation policies across clean and realistic starting conditions.")

    slide = blank_slide(prs)
    add_title(slide, "Algorithms Compared")
    rows = [["Family", "Policy", "What it knows", "Main cost"],
            ["Stateless", "Random", "one random bin", "very low"],
            ["Stateless", "Power-k", "k sampled bin loads", "grows with k"],
            ["Stateless", "Absolute minimum", "all bin loads", "high scan cost"],
            ["Stateful", "Round-robin", "next-bin pointer", "very low state"],
            ["Stateful", "Heap size s", "partial/global heap state", "state + heap updates"]]
    add_table(slide, 0.65, 1.45, 12.0, 4.1, rows)
    add_takeaway(slide, "The core comparison is stateless sampling versus maintained state versus full global knowledge.")

    slide = blank_slide(prs)
    add_title(slide, "Theory Anchor: Why Power-2 Matters")
    add_bullets(slide, 1.1, 1.65, 6.2, 3.1,
                ["Random allocation has a large maximum-load tail.",
                 "Power of Two Choices gives a dramatic theoretical improvement.",
                 "Additional choices help, but with diminishing returns.",
                 "The simulation checks whether this holds with cost, weights, and initialization."], 20)
    add_callout(slide, 8.0, 2.1, 3.8, 1.35,
                ["Hypothesis", "Small local information captures most of the benefit."], "blue")
    add_takeaway(slide, "Theory motivates Power-2 as the practical randomized baseline.")

    slide = blank_slide(prs)
    add_title(slide, "Clean Baseline: Power-2 Captures Most of the Gain")
    add_plot(slide, "plots/08_main_comparison_unweighted_empty_max_load.png")
    add_callout(slide, 9.55, 1.6, 2.9, 1.65,
                ["Max load", "n=16: Random 666.1 -> Power-2 626.4",
                 "n=32: Random 345.3 -> Power-2 313.9"], "blue")
    add_callout(slide, 9.55, 3.65, 2.9, 1.0,
                ["Beyond Power-2", "Absolute minimum adds only a small improvement."], "gray")
    add_takeaway(slide, "Power-2 nearly closes the gap to the ideal policy in the clean unweighted case.")

    slide = blank_slide(prs)
    add_title(slide, "Diminishing Returns of k")
    add_plot(slide, "plots/02_diminishing_returns_max_load.png")
    add_callout(slide, 9.55, 1.55, 2.9, 1.8,
                ["n=16", "Power-2/3/4 max load: 626.4 / 626.0 / 625.9",
                 "Cost per ball: 9 / 14 / 19"], "blue")
    add_callout(slide, 9.55, 3.8, 2.9, 1.0,
                ["n=32", "Power-2/3/4: 313.9 / 313.6 / 313.0"], "blue")
    add_takeaway(slide, "The big improvement is k=1 to k=2. Larger k mostly increases cost.")

    slide = blank_slide(prs)
    add_title(slide, "Cost-Quality Tradeoff")
    add_plot(slide, "plots/09_tradeoff_unweighted_empty_max_load.png", w=8.7, h=3.75)
    add_bullets(slide, 9.65, 1.55, 2.8, 3.7,
                ["Lower max load is better.",
                 "Lower cost per ball is better.",
                 "Absolute minimum is close to optimal but expensive.",
                 "Round-robin looks excellent here, but that success is fragile."], 15)
    add_takeaway(slide, "Power-2 is the knee of the curve: most of the quality for a small cost.")

    slide = blank_slide(prs)
    add_title(slide, "Weighted Balls Break Count Fairness")
    add_plot(slide, "plots/04_weighted_round_robin_break_max_load.png")
    add_callout(slide, 9.55, 1.55, 2.9, 1.2,
                ["Weighted-empty n=32", "Round-robin 1812.6", "Power-2 1731.6"], "green")
    add_bullets(slide, 9.65, 3.25, 2.75, 1.5,
                ["Fair by count can still be bad by load.",
                 "Power-2 reads current load and adapts."], 15)
    add_takeaway(slide, "Round-robin balances counts, not load. Weighted balls expose that weakness.", color="green")

    slide = blank_slide(prs)
    add_title(slide, "Stateful Heap: When Maintained State Helps")
    add_plot(slide, "plots/05_heap_sweep_n32.png")
    add_callout(slide, 9.55, 1.55, 2.9, 1.55,
                ["Full heap", "Matches absolute-min quality around 1724.7",
                 "Cost around 38.32 instead of 159"], "orange")
    add_bullets(slide, 9.65, 3.6, 2.75, 1.7,
                ["Partial heap tracks s bins.",
                 "Full heap tracks all bins.",
                 "More state improves quality and update cost."], 15)
    add_takeaway(slide, "A maintained heap can approach global-min quality without scanning all bins.", color="orange")

    slide = blank_slide(prs)
    add_title(slide, "Initialization Tests Robustness")
    add_plot(slide, "plots/08_main_comparison_unweighted_initialized_max_load.png")
    add_callout(slide, 9.55, 1.55, 2.9, 1.55,
                ["Near reference", "n=16: Power-2 679.3 vs ref 678.0",
                 "n=32: Power-2 365.4 vs ref 364.0"], "blue")
    add_bullets(slide, 9.65, 3.6, 2.75, 1.45,
                ["Real systems rarely start from zero.",
                 "Round-robin can preserve initial imbalance."], 15)
    add_takeaway(slide, "Starting from non-empty bins exposes whether a policy corrects imbalance.")

    slide = blank_slide(prs)
    add_title(slide, "Final Realistic Tradeoff")
    add_plot(slide, "plots/09_tradeoff_weighted_initialized_max_load.png", w=8.7, h=3.75)
    add_bullets(slide, 9.65, 1.45, 2.85, 3.95,
                ["Random is too weak.",
                 "Absolute minimum is expensive.",
                 "Round-robin is fragile.",
                 "Power-2 is simple, stateless, and robust.",
                 "Heap/full heap can improve quality when state cost is justified."], 14.5)
    add_takeaway(slide, "Under weighted + initialized conditions, Power-2 is the strongest simple default.")

    slide = blank_slide(prs)
    add_title(slide, "C++ Implementation Overview")
    add_pipeline(slide)
    add_bullets(slide, 0.85, 1.45, 5.5, 1.25,
                ["src/main.cpp builds scenarios, runs them, and exports results/simulation_results.csv.",
                 "ExperimentRunner creates all experiment scenarios."], 15)
    add_bullets(slide, 0.85, 4.35, 5.8, 1.35,
                ["PowerKSimulator: Random, Power-k, Absolute minimum.",
                 "StatefulRoundRobinSimulator: round-robin.",
                 "HeapSizeSPowerOfKSimulator: heap size s / full heap."], 15)
    add_bullets(slide, 7.15, 4.35, 5.2, 1.0,
                ["scripts/plot_results.py reads the CSV and creates the plots.",
                 "The deck uses repository outputs only."], 15)
    add_takeaway(slide, "The project is structured as a repeatable performance evaluation pipeline.")

    slide = blank_slide(prs)
    add_title(slide, "Final Takeaway")
    add_textbox(slide, Inches(1.0), Inches(1.65), Inches(11.1), Inches(1.15),
                "Two choices are the best default.", 30, True, "blue",
                PP_ALIGN.CENTER)
    add_bullets(slide, 2.0, 3.0, 9.4, 2.0,
                ["Most of the balance benefit with little decision cost.",
                 "Good robustness across weights and initialized bins.",
                 "Full global knowledge is rarely worth the price.",
                 "Stateful heap designs can improve quality further when memory/update cost is justified."],
                20)
    add_takeaway(slide, "Engineering decision: optimize load-balance quality per unit of decision and state cost.")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build_deck()
